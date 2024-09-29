import json
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
import os
import logging
from dotenv import load_dotenv

load_dotenv()



# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Unsplash API access key - replace with your own
UNSPLASH_ACCESS_KEY = os.getenv("UNSPLASH_ACCESS_KEY")

def download_image(query):
    search_url = f"https://api.unsplash.com/search/photos"
    params = {
        "query": query,
        "per_page": 1,
        "client_id": UNSPLASH_ACCESS_KEY
    }

    try:
        response = requests.get(search_url, params=params, timeout=10)
        response.raise_for_status()
        data = response.json()

        if data["results"]:
            img_url = data["results"][0]["urls"]["regular"]
            img_response = requests.get(img_url, timeout=5)
            img_response.raise_for_status()
            img_data = img_response.content

            img_path = f"images/{query.replace(' ', '_')}.jpg"
            with open(img_path, 'wb') as handler:
                handler.write(img_data)
            logger.info(f"Successfully downloaded image for query: {query}")
            return img_path
        else:
            logger.warning(f"No images found for query: {query}")
    except Exception as e:
        logger.error(f"Error during image search for query {query}: {str(e)}")

    return None

def create_presentation(slides_data):
    prs = Presentation()
    layout = prs.slide_layouts[5]  # Layout with a title and content

    if not os.path.exists('images'):
        os.makedirs('./images')

    for slide_data in slides_data:
        slide = prs.slides.add_slide(layout)
        
        # Set the slide title
        title = slide.shapes.title
        title.text = slide_data['title']

        # Download and add the image
        image_path = download_image(slide_data['image'])
        if image_path:
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(8)
            height = Inches(4.5)
            slide.shapes.add_picture(image_path, left, top, width, height)
        else:
            logger.error(f"Failed to download image for query: {slide_data['image']}")

    prs.save("Ansible_Presentation.pptx")
